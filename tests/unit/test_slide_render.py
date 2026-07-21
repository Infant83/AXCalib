from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from axcalib.ingest import EmbeddedImagePptxRenderer, SlideRenderError

ROOT = Path(__file__).resolve().parents[2]
SOURCE = ROOT / "tests" / "sources" / "oled_qc_project_outline.pptx"


def test_restricted_renderer_covers_actual_pptx_and_is_reproducible(
    tmp_path: Path,
) -> None:
    renderer = EmbeddedImagePptxRenderer()

    first = renderer.render(SOURCE, tmp_path / "first")
    second = renderer.render(SOURCE, tmp_path / "second")

    assert first.slide_count == 16
    assert first.rendered_slide_count == 16
    assert first.visual_slide_count == 13
    assert first.blank_slide_count == 3
    assert {item.slide_number for item in first.artifacts if not item.visual_content_present} == {
        6,
        15,
        16,
    }
    assert {(item.width_px, item.height_px) for item in first.artifacts} == {(1672, 941)}
    assert first.canonical_sha256 == second.canonical_sha256
    assert [item.image_sha256 for item in first.artifacts] == [
        item.image_sha256 for item in second.artifacts
    ]
    assert all((tmp_path / "first" / item.artifact_uri).is_file() for item in first.artifacts)
    assert (tmp_path / "first" / first.manifest_uri).is_file()


def test_restricted_renderer_fails_closed_for_composed_slide(tmp_path: Path) -> None:
    source = tmp_path / "composed.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "text"
    presentation.save(str(source))

    with pytest.raises(SlideRenderError, match="composed content"):
        EmbeddedImagePptxRenderer().render(source, tmp_path / "render")
