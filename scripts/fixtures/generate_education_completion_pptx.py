"""Generate the clearly labeled synthetic completion PPTX fixture and sidecar."""

# pyright: reportMissingImports=false

from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDES = (
    (
        "Synthetic 완료평가 보고서",
        (
            "과제: OLED 분자 역설계 프로젝트",
            "결과: 교육 lifecycle과 AXCalib 두 Gate를 검증한 테스트 전용 완료자료",
            "산출물: dossier, 등록·완료 리포트, 감사 로그, 재현 스크립트",
            "주의: 실제 연구성과나 실제 학습자 평가가 아님",
        ),
        ("deliverable", "result"),
    ),
    (
        "등록 baseline과 수행 변경",
        (
            "등록 목표: SELFIES·surrogate·QUBO를 연결한 역설계 검증 계획 수립",
            "변경: 실제 양자 하드웨어 실험 대신 deterministic offline workflow 검증으로 범위 조정",
            "변경 승인: synthetic 교육 예제 범위에 한정한 관리자 테스트 입력",
        ),
        ("change", "result"),
    ),
    (
        "KPI 계획 대비 관측값",
        (
            "KPI 1 · 두 HITL 승인요청 기록: target 2건, observed 2건",
            "KPI 2 · 단계 누출: target 0건, observed 0건",
            "KPI 3 · criterion locator 보존: target 100%, observed 100%",
            "측정기간: 단일 offline integration run",
        ),
        ("result", "quantitative_target", "kpi_plan"),
    ),
    (
        "수행·재현 증거",
        (
            "코드 버전과 program version을 실행기록에 고정",
            "PPTX와 sidecar SHA-256 hash를 dossier artifact로 기록",
            "테스트 로그와 JSON·Markdown 리포트를 workspace에 생성",
            "재현 절차: 문서화된 Python library example을 동일 fixture로 실행",
        ),
        ("deliverable", "result", "reproducibility"),
    ),
    (
        "남은 리스크와 한계",
        (
            "리스크: 실제 embedding·Qdrant·on-prem Qwen 품질은 검증하지 않음",
            "한계: 수동 reviewed sidecar는 OCR/VLM gold label이 아님",
            "후속계획: 평가 책임자 rubric 승인 후 비식별 gold dataset으로 별도 검증",
        ),
        ("risk", "limitation"),
    ),
    (
        "완료 증거 manifest",
        (
            "산출물 reference와 content hash를 dossier에서 추적",
            "결과 report는 Agent 제안과 관리자 결정을 분리",
            "재현 테스트는 외부 endpoint 없이 실행",
            "본 문서는 SYNTHETIC EXAMPLE이며 공식 인증서가 아님",
        ),
        ("deliverable", "result", "reproducibility"),
    ),
)


def build_deck(target: Path) -> None:
    """Create a local-only deck with standard OOXML text nodes."""

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide_height = presentation.slide_height
    if slide_height is None:
        raise RuntimeError("presentation slide height was not initialized")
    for index, (title, bullets, _) in enumerate(SLIDES, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(247, 249, 252)

        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(0.18),
            slide_height,
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(22, 74, 154)
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(
            Inches(0.65), Inches(0.55), Inches(12.0), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(17, 42, 81)

        body_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(1.65), Inches(11.5), Inches(4.7)
        )
        body = body_box.text_frame
        body.clear()
        body.word_wrap = True
        for bullet_index, text in enumerate(bullets):
            item = body.paragraphs[0] if bullet_index == 0 else body.add_paragraph()
            item.text = text
            item.level = 0
            item.font.size = Pt(20)
            item.font.color.rgb = RGBColor(35, 48, 68)
            item.space_after = Pt(16)

        footer = slide.shapes.add_textbox(
            Inches(0.65), Inches(6.85), Inches(12.0), Inches(0.35)
        )
        footer_p = footer.text_frame.paragraphs[0]
        footer_p.text = f"SYNTHETIC EXAMPLE · AXCalib education lifecycle · {index}/{len(SLIDES)}"
        footer_p.alignment = PP_ALIGN.RIGHT
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(100, 110, 125)

    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))


def write_sidecar(source: Path, target: Path) -> None:
    digest = hashlib.sha256(source.read_bytes()).hexdigest()
    payload = {
        "schema_version": "axcalib.pptx-sidecar/v1",
        "source_sha256": digest,
        "annotation_status": "synthetic_fixture_author_annotation",
        "slides": [
            {
                "slide": index,
                "summary": " ".join(bullets),
                "tags": list(tags),
            }
            for index, (_, bullets, tags) in enumerate(SLIDES, start=1)
        ],
    }
    target.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=(
            Path(__file__).resolve().parents[2]
            / "fixtures"
            / "synthetic"
            / "education_project_lifecycle"
        ),
    )
    args = parser.parse_args()
    pptx = args.output_dir / "completion_report.synthetic.pptx"
    sidecar = args.output_dir / "completion_report.synthetic.axcalib.json"
    build_deck(pptx)
    write_sidecar(pptx, sidecar)
    print(pptx.resolve())
    print(sidecar.resolve())
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
